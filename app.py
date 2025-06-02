from flask import Flask, request, jsonify
import os

app = Flask(__name__)

# ⏳ 延遲初始化：首次 API 請求才載入模型與文件
qa = None

@app.route("/ask", methods=["POST"])
def ask():
    global qa
    data = request.get_json()
    query = data.get("question")

    if not query:
        return jsonify({"error": "缺少問題內容"}), 400

    # 首次請求才初始化模型與向量資料庫
    if qa is None:
        try:
            from langchain.embeddings import HuggingFaceEmbeddings
            from langchain.vectorstores import FAISS
            from langchain.text_splitter import RecursiveCharacterTextSplitter
            from langchain.llms import HuggingFaceHub
            from langchain.chains import RetrievalQA
            import pptx
            import docx

            def extract_text_from_pptx(path):
                prs = pptx.Presentation(path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)

            def extract_text_from_docx(path):
                doc = docx.Document(path)
                return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

            def load_documents():
                texts = []
                for fname in os.listdir("docs"):
                    fpath = os.path.join("docs", fname)
                    if fname.endswith(".pptx"):
                        texts.append(extract_text_from_pptx(fpath))
                    elif fname.endswith(".docx"):
                        texts.append(extract_text_from_docx(fpath))
                return "\n".join(texts)

            print("✅ 正在初始化模型與文件...")
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
            documents = load_documents()
            texts = text_splitter.split_text(documents)
            embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
            db = FAISS.from_texts(texts, embeddings)
            retriever = db.as_retriever()

            llm = HuggingFaceHub(
                repo_id="google/flan-t5-base",
                model_kwargs={"temperature": 0.5, "max_new_tokens": 256}
            )
            qa = RetrievalQA.from_chain_type(llm=llm, retriever=retriever)
            print("✅ 初始化完成")
        except Exception as e:
            return jsonify({"error": f"初始化失敗：{str(e)}"}), 500

    try:
        result = qa.run(query)
        return jsonify({"answer": result})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# 🟢 為 Render 設定正確的 host/port
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    print(f"🚀 App running on port {port}")
    app.run(host="0.0.0.0", port=port)
