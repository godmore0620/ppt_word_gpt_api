# ✅ Flask API：吃你的 PPT + Word 檔案來回應提問
# 使用：FAISS + HuggingFace 模型 + flask + langchain

from flask import Flask, request, jsonify
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.llms import HuggingFaceHub
from langchain.chains import RetrievalQA
import os
import tempfile
import pptx
import docx

app = Flask(__name__)

# ✅ 把 PPT 和 Word 檔案轉成文字
def extract_text_from_pptx(path):
    prs = pptx.Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

# ✅ 載入並分段文字資料（你要的資料庫）
def load_documents():
    texts = []
    for fname in os.listdir("docs"):
        fpath = os.path.join("docs", fname)
        if fname.endswith(".pptx"):
            texts.append(extract_text_from_pptx(fpath))
        elif fname.endswith(".docx"):
            texts.append(extract_text_from_docx(fpath))
    return "\n".join(texts)

# ✅ 建立向量資料庫
print("\u2705 正在處理文件...")
text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
documents = load_documents()
texts = text_splitter.split_text(documents)
embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
db = FAISS.from_texts(texts, embeddings)
retriever = db.as_retriever()

# ✅ HuggingFace 模型（可改成 GPT3.5）
llm = HuggingFaceHub(repo_id="tiiuae/falcon-7b-instruct", model_kwargs={"temperature":0.5, "max_new_tokens":256})
qa = RetrievalQA.from_chain_type(llm=llm, retriever=retriever)

@app.route("/ask", methods=["POST"])
def ask():
    data = request.get_json()
    query = data.get("question")
    if not query:
        return jsonify({"error": "缺少問題內容。"}), 400
    try:
        result = qa.run(query)
        return jsonify({"answer": result})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(debug=True)
